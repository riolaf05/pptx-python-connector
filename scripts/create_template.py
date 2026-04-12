"""
Script per generare il file PPT template con placeholder.

I placeholder sono stringhe nel formato {{nome_campo}} inserite
nelle caselle di testo delle slide. La Lambda li sostituisce con
i valori ricevuti nel JSON di input.

Esegui:
    pip install python-pptx
    python scripts/create_template.py

Il file viene salvato in examples/report_template.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import os

OUTPUT_PATH = os.path.join(os.path.dirname(__file__), "..", "examples", "report_template.pptx")


def add_textbox(slide, left, top, width, height, text, font_size=14,
                bold=False, color=RGBColor(0x33, 0x33, 0x33),
                alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    """Helper: aggiunge una textbox alla slide."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                      Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_shape_rect(slide, left, top, width, height,
                   fill_color=RGBColor(0x00, 0x56, 0x9E)):
    """Helper: aggiunge un rettangolo colorato come sfondo."""
    from pptx.enum.shapes import MSO_SHAPE
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def create_template():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # =========================================================
    # SLIDE 1 - Copertina
    # =========================================================
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # layout vuoto

    # Sfondo colorato header
    add_shape_rect(slide1, 0, 0, 13.333, 4.2, RGBColor(0x00, 0x56, 0x9E))

    # Titolo
    add_textbox(slide1, 0.8, 1.0, 11.5, 1.2,
                "{{titolo_report}}",
                font_size=40, bold=True,
                color=RGBColor(0xFF, 0xFF, 0xFF),
                alignment=PP_ALIGN.LEFT)

    # Sottotitolo
    add_textbox(slide1, 0.8, 2.4, 11.5, 0.8,
                "{{sottotitolo}}",
                font_size=22,
                color=RGBColor(0xCC, 0xDD, 0xEE),
                alignment=PP_ALIGN.LEFT)

    # Autore e data
    add_textbox(slide1, 0.8, 4.8, 5, 0.5,
                "Autore: {{autore}}",
                font_size=16, color=RGBColor(0x55, 0x55, 0x55))

    add_textbox(slide1, 0.8, 5.4, 5, 0.5,
                "Data: {{data_report}}",
                font_size=16, color=RGBColor(0x55, 0x55, 0x55))

    add_textbox(slide1, 0.8, 6.0, 5, 0.5,
                "Dipartimento: {{dipartimento}}",
                font_size=16, color=RGBColor(0x55, 0x55, 0x55))

    # =========================================================
    # SLIDE 2 - KPI Overview
    # =========================================================
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])

    # Header bar
    add_shape_rect(slide2, 0, 0, 13.333, 1.2, RGBColor(0x00, 0x56, 0x9E))
    add_textbox(slide2, 0.5, 0.15, 12, 0.9,
                "KPI Overview",
                font_size=30, bold=True,
                color=RGBColor(0xFF, 0xFF, 0xFF))

    # KPI Cards - riga 1
    kpi_data = [
        ("Fatturato Totale", "{{fatturato_totale}}", 0.5),
        ("Variazione %", "{{variazione_percentuale}}", 3.7),
        ("Clienti Totali", "{{numero_clienti}}", 6.9),
        ("Nuovi Clienti", "{{nuovi_clienti}}", 10.1),
    ]
    for label, value, left in kpi_data:
        add_shape_rect(slide2, left, 1.8, 2.8, 1.8, RGBColor(0xF0, 0xF4, 0xF8))
        add_textbox(slide2, left + 0.2, 1.9, 2.4, 0.5,
                    label, font_size=13, color=RGBColor(0x66, 0x66, 0x66))
        add_textbox(slide2, left + 0.2, 2.5, 2.4, 0.8,
                    value, font_size=26, bold=True,
                    color=RGBColor(0x00, 0x56, 0x9E))

    # Obiettivo
    add_shape_rect(slide2, 0.5, 4.2, 6.0, 1.5, RGBColor(0xE8, 0xF5, 0xE9))
    add_textbox(slide2, 0.7, 4.3, 5.5, 0.5,
                "Obiettivo Trimestrale", font_size=14,
                color=RGBColor(0x66, 0x66, 0x66))
    add_textbox(slide2, 0.7, 4.8, 2.5, 0.7,
                "{{obiettivo_trimestrale}}", font_size=24, bold=True,
                color=RGBColor(0x2E, 0x7D, 0x32))
    add_textbox(slide2, 3.5, 4.8, 2.5, 0.7,
                "{{stato_obiettivo}}", font_size=24, bold=True,
                color=RGBColor(0x2E, 0x7D, 0x32),
                alignment=PP_ALIGN.RIGHT)

    # =========================================================
    # SLIDE 3 - Top Prodotti
    # =========================================================
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])

    add_shape_rect(slide3, 0, 0, 13.333, 1.2, RGBColor(0x00, 0x56, 0x9E))
    add_textbox(slide3, 0.5, 0.15, 12, 0.9,
                "Top Prodotti",
                font_size=30, bold=True,
                color=RGBColor(0xFF, 0xFF, 0xFF))

    products = [
        ("1", "{{prodotto_top_1}}", "{{prodotto_top_1_valore}}"),
        ("2", "{{prodotto_top_2}}", "{{prodotto_top_2_valore}}"),
        ("3", "{{prodotto_top_3}}", "{{prodotto_top_3_valore}}"),
    ]
    for i, (rank, name, value) in enumerate(products):
        y = 1.8 + i * 1.6
        add_shape_rect(slide3, 0.5, y, 12.3, 1.3, RGBColor(0xF5, 0xF5, 0xF5))

        # Rank badge
        add_shape_rect(slide3, 0.7, y + 0.2, 0.8, 0.8, RGBColor(0x00, 0x56, 0x9E))
        add_textbox(slide3, 0.7, y + 0.25, 0.8, 0.8,
                    f"#{rank}", font_size=22, bold=True,
                    color=RGBColor(0xFF, 0xFF, 0xFF),
                    alignment=PP_ALIGN.CENTER)

        # Product name
        add_textbox(slide3, 1.8, y + 0.3, 6, 0.7,
                    name, font_size=20, bold=True,
                    color=RGBColor(0x33, 0x33, 0x33))

        # Product value
        add_textbox(slide3, 9.0, y + 0.3, 3.5, 0.7,
                    value, font_size=22, bold=True,
                    color=RGBColor(0x00, 0x56, 0x9E),
                    alignment=PP_ALIGN.RIGHT)

    # =========================================================
    # SLIDE 4 - Note e Osservazioni
    # =========================================================
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])

    add_shape_rect(slide4, 0, 0, 13.333, 1.2, RGBColor(0x00, 0x56, 0x9E))
    add_textbox(slide4, 0.5, 0.15, 12, 0.9,
                "Note e Osservazioni",
                font_size=30, bold=True,
                color=RGBColor(0xFF, 0xFF, 0xFF))

    notes = [
        ("1", "{{nota_1}}"),
        ("2", "{{nota_2}}"),
        ("3", "{{nota_3}}"),
    ]
    for i, (num, text) in enumerate(notes):
        y = 1.8 + i * 1.5
        add_shape_rect(slide4, 0.5, y, 12.3, 1.2, RGBColor(0xFD, 0xF5, 0xE6))
        add_textbox(slide4, 0.8, y + 0.15, 0.5, 0.8,
                    num + ".", font_size=20, bold=True,
                    color=RGBColor(0xE6, 0x8A, 0x00))
        add_textbox(slide4, 1.4, y + 0.2, 11, 0.8,
                    text, font_size=16,
                    color=RGBColor(0x44, 0x44, 0x44))

    # =========================================================
    # SLIDE 5 - Contatti
    # =========================================================
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])

    add_shape_rect(slide5, 0, 0, 13.333, 7.5, RGBColor(0x00, 0x56, 0x9E))

    add_textbox(slide5, 0.5, 1.5, 12, 1.0,
                "Contatti",
                font_size=36, bold=True,
                color=RGBColor(0xFF, 0xFF, 0xFF),
                alignment=PP_ALIGN.CENTER)

    add_textbox(slide5, 0.5, 3.0, 12, 0.6,
                "{{nome_contatto}}",
                font_size=24, bold=True,
                color=RGBColor(0xFF, 0xFF, 0xFF),
                alignment=PP_ALIGN.CENTER)

    add_textbox(slide5, 0.5, 3.8, 12, 0.6,
                "{{email_contatto}}",
                font_size=18,
                color=RGBColor(0xCC, 0xDD, 0xEE),
                alignment=PP_ALIGN.CENTER)

    add_textbox(slide5, 0.5, 4.5, 12, 0.6,
                "{{telefono_contatto}}",
                font_size=18,
                color=RGBColor(0xCC, 0xDD, 0xEE),
                alignment=PP_ALIGN.CENTER)

    # Salva
    out = os.path.abspath(OUTPUT_PATH)
    prs.save(out)
    print(f"Template salvato in: {out}")


if __name__ == "__main__":
    create_template()
