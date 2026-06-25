"""
Script per generare examples/report_gtm_template.pptx.

Carica il template originale GTM Foundation Report e sostituisce il contenuto
variabile con placeholder {{chiave}} compatibili con la Lambda gtm_handler.

Utilizzo:
    python scripts/create_gtm_template.py [percorso_input]

Se il percorso non è specificato, usa DEFAULT_INPUT_PATH.
"""

import os
import sys
from pptx import Presentation

DEFAULT_INPUT_PATH = r"C:\Users\r.laface\Desktop\VA Agents\GTM Foundation Report template.pptx"
OUTPUT_PATH = os.path.join(os.path.dirname(__file__), "..", "examples", "report_gtm_template.pptx")

# Regole per sostituzione a livello di run (primo paragrafo dello shape).
# Usato per Slide 1 dove ogni cella è un run distinto nella stessa paragraph.
# {(slide_idx, shape_name, occurrence): {run_idx: new_text}}
RUN_RULES = {
    (0, "Title 3", 0): {
        0: "{{customer_name}}",
    },
    (0, "Content Placeholder 4", 0): {
        1: "Date: {{document_date}}",
        2: "Version: {{document_version}}",
    },
}

# Regole per sostituzione a livello di paragrafo (primo run del paragrafo).
# {(slide_idx, shape_name, occurrence): {para_idx: new_text}}
PARA_RULES = {
    # --- Slide 5: ICP ---
    (4, "Cell0Content", 0): {0: "{{icp_industry}}"},
    (4, "Cell1Content", 0): {0: "{{icp_geography}}"},
    (4, "Cell2Content", 0): {0: "{{icp_company_size}}"},
    (4, "Cell0Content", 1): {0: "{{icp_buying_group}}"},
    (4, "Cell1Content", 1): {0: "{{icp_maturity_level}}"},
    (4, "Cell2Content", 1): {0: "{{icp_custom_field}}"},

    # --- Slide 7: Persona ---
    (6, "RoleName", 0): {0: "{{persona_role_name}}"},
    (6, "BasicInfo", 0): {
        0: "Name: {{persona_name}}",
        1: "Role: {{persona_role}}",
        2: "Age: {{persona_age}}",
        3: "Education: {{persona_education}}",
        4: "Location: {{persona_location}}",
        5: "Industry: {{persona_industry}}",
    },
    (6, "RespContent", 0): {
        0: "{{persona_resp_1}}",
        1: "{{persona_resp_2}}",
        2: "{{persona_resp_3}}",
    },
    (6, "CardHead0", 0): {0: "{{persona_obj1_title}}"},
    (6, "CardBody0", 0): {
        0: "Description: {{persona_obj1_desc}}",
        1: "KPI 1: {{persona_obj1_kpi1}}",
        2: "KPI 2: {{persona_obj1_kpi2}}",
        3: "KPI 3: {{persona_obj1_kpi3}}",
        4: "Pain Point 1: {{persona_obj1_pain1}}",
        5: "Pain Point 2: {{persona_obj1_pain2}}",
        6: "Winning Trigger: {{persona_obj1_trigger}}",
    },
    (6, "CardHead0", 1): {0: "{{persona_obj2_title}}"},
    (6, "CardBody0", 1): {
        0: "Description: {{persona_obj2_desc}}",
        1: "KPI 1: {{persona_obj2_kpi1}}",
        2: "KPI 2: {{persona_obj2_kpi2}}",
        3: "KPI 3: {{persona_obj2_kpi3}}",
        4: "Pain Point 1: {{persona_obj2_pain1}}",
        5: "Pain Point 2: {{persona_obj2_pain2}}",
        6: "Winning Trigger: {{persona_obj2_trigger}}",
    },
    (6, "CardHead0", 2): {0: "{{persona_obj3_title}}"},
    (6, "CardBody0", 2): {
        0: "Description: {{persona_obj3_desc}}",
        1: "KPI 1: {{persona_obj3_kpi1}}",
        2: "KPI 2: {{persona_obj3_kpi2}}",
        3: "KPI 3: {{persona_obj3_kpi3}}",
        4: "Pain Point 1: {{persona_obj3_pain1}}",
        5: "Pain Point 2: {{persona_obj3_pain2}}",
        6: "Winning Trigger: {{persona_obj3_trigger}}",
    },
    (6, "CardHead0", 3): {0: "{{persona_obj4_title}}"},
    (6, "CardBody0", 3): {
        0: "Description: {{persona_obj4_desc}}",
        1: "KPI 1: {{persona_obj4_kpi1}}",
        2: "KPI 2: {{persona_obj4_kpi2}}",
        3: "KPI 3: {{persona_obj4_kpi3}}",
        4: "Pain Point 1: {{persona_obj4_pain1}}",
        5: "Pain Point 2: {{persona_obj4_pain2}}",
        6: "Winning Trigger: {{persona_obj4_trigger}}",
    },
    (6, "CardBody0", 4): {0: "{{persona_channels}}"},

    # --- Slide 10: Value Proposition ---
    (9, "A0", 0): {0: "{{vp_what_is}}"},
    (9, "A1", 0): {0: "{{vp_what_does}}"},
    (9, "A2", 0): {0: "{{vp_how_works}}"},
    (9, "A3", 0): {0: "{{vp_value}}"},

    # --- Slide 11: Slogans ---
    (10, "SlContent0", 0): {0: "{{slogan_1}}"},
    (10, "SlContent1", 0): {0: "{{slogan_2}}"},
    (10, "SlContent2", 0): {0: "{{slogan_3}}"},

    # --- Slide 13: Customer Mental Model ---
    # para[0] = "CUSTOMER MENTAL MODEL" (titolo fisso), para[1] = nome persona
    (12, "TitleText", 0): {1: "{{cmm_persona}}"},
    (12, "ObjA", 0):  {0: "{{cmm_objective}}"},
    (12, "DF1", 0):   {0: "{{cmm_driving_factor_1}}"},
    (12, "DF2", 0):   {0: "{{cmm_driving_factor_2}}"},
    (12, "PP1", 0):   {0: "{{cmm_pain_point_1}}"},
    (12, "PP2", 0):   {0: "{{cmm_pain_point_2}}"},
    (12, "UC1", 0):   {0: "{{cmm_use_case_1}}"},
    (12, "UC2", 0):   {0: "{{cmm_use_case_2}}"},

    # --- Slide 15: USPs ---
    (14, "USP0HdrTxt", 0): {0: "{{usp_1_title}}"},
    (14, "USP0Body", 0): {0: "{{usp_1_body_1}}", 1: "{{usp_1_body_2}}", 2: "{{usp_1_body_3}}"},
    (14, "USP1HdrTxt", 0): {0: "{{usp_2_title}}"},
    (14, "USP1Body", 0): {0: "{{usp_2_body_1}}", 1: "{{usp_2_body_2}}", 2: "{{usp_2_body_3}}"},
    (14, "USP2HdrTxt", 0): {0: "{{usp_3_title}}"},
    (14, "USP2Body", 0): {0: "{{usp_3_body_1}}", 1: "{{usp_3_body_2}}", 2: "{{usp_3_body_3}}"},

    # --- Slide 17: Commercial Insight ---
    (16, "Node0Val", 0): {0: "{{ci_objective}}"},
    (16, "Node1Val", 0): {0: "{{ci_pain_point}}"},
    (16, "Node2Val", 0): {0: "{{ci_use_case}}"},
    (16, "Node3Val", 0): {0: "{{ci_usp}}"},
    (16, "InsightTitle", 0): {0: "{{ci_title}}"},
    (16, "InsightBody", 0): {
        0: "{{ci_narrative_1}}",
        1: "{{ci_narrative_2}}",
        2: "{{ci_narrative_3}}",
    },
}


def _set_run(para, run_idx: int, new_text: str) -> None:
    if run_idx < len(para.runs):
        para.runs[run_idx].text = new_text


def _set_para_first_run(para, new_text: str) -> None:
    """Sostituisce il testo nel paragrafo usando solo il primo run."""
    if para.runs:
        para.runs[0].text = new_text
        for run in para.runs[1:]:
            run.text = ""
    else:
        para.text = new_text


def apply_rules(prs: Presentation) -> None:
    for slide_idx, slide in enumerate(prs.slides):
        name_count: dict[str, int] = {}
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            n = shape.name
            occ = name_count.get(n, 0)
            name_count[n] = occ + 1
            key = (slide_idx, n, occ)

            if key in RUN_RULES:
                first_para = shape.text_frame.paragraphs[0]
                for run_idx, new_text in RUN_RULES[key].items():
                    _set_run(first_para, run_idx, new_text)

            if key in PARA_RULES:
                tf = shape.text_frame
                for para_idx, new_text in PARA_RULES[key].items():
                    if para_idx < len(tf.paragraphs):
                        _set_para_first_run(tf.paragraphs[para_idx], new_text)


def create_gtm_template(input_path: str | None = None) -> None:
    src = input_path or DEFAULT_INPUT_PATH
    prs = Presentation(src)
    apply_rules(prs)
    out = os.path.abspath(OUTPUT_PATH)
    os.makedirs(os.path.dirname(out), exist_ok=True)
    prs.save(out)
    print(f"GTM template salvato in: {out}")


if __name__ == "__main__":
    inp = sys.argv[1] if len(sys.argv) > 1 else None
    create_gtm_template(inp)
