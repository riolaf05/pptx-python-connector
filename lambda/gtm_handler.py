"""
AWS Lambda handler – compila il GTM Foundation Report template PPTX
sostituendo i placeholder {{chiave}} con i valori ricevuti nel JSON di input,
salva il risultato su S3 e restituisce un pre-signed URL per il download.

Richiede l'header x-api-key per l'autenticazione.

Formato di input atteso:
{
  "customer_name": "...",
  "document_date": "DD.MM.YYYY",
  "document_version": "1.0",
  "icp": { "industry": "...", "geography": "...", ... },
  "persona": { "role_name": "...", "name": "...", "objective_1": {...}, ... },
  "value_proposition": { "what_is": "...", ... },
  "slogans": ["...", "...", "..."],
  "customer_mental_model": { "persona": "...", "objective": "...", ... },
  "usps": [{ "title": "...", "body": [...] }, ...],
  "commercial_insight": { "title": "...", "narratives": [...], ... }
}

In alternativa, formato diretto con chiave "placeholders":
{
  "template_s3_key": "templates/report_gtm_template.pptx",
  "placeholders": { "customer_name": "...", "icp_industry": "...", ... }
}
"""

import json
import os
import re
import tempfile
import traceback

import boto3
from pptx import Presentation

S3_BUCKET                = os.environ["S3_BUCKET"]
API_KEY                  = os.environ.get("API_KEY", "")
GTM_TEMPLATE_KEY         = os.environ.get("GTM_TEMPLATE_KEY", "templates/report_gtm_template.pptx")
PRESIGNED_URL_EXPIRATION = int(os.environ.get("PRESIGNED_URL_EXPIRATION", "3600"))
ALERT_EMAIL              = os.environ.get("ALERT_EMAIL", "")
ALERT_FROM_EMAIL         = os.environ.get("ALERT_FROM_EMAIL", "")

s3  = boto3.client("s3")
ses = boto3.client("ses")

PLACEHOLDER_RE = re.compile(r"\{\{(\w+)\}\}")


def check_api_key(event: dict) -> bool:
    """Verifica che l'header x-api-key sia presente e corretto."""
    if not API_KEY:
        return True
    headers = event.get("headers") or {}
    provided = headers.get("x-api-key") or headers.get("X-Api-Key") or ""
    return provided == API_KEY


def _expand_list(items: list, prefix: str, max_n: int) -> dict:
    return {
        f"{prefix}_{i}": items[i - 1] if i <= len(items) else ""
        for i in range(1, max_n + 1)
    }


def preprocess_gtm(body: dict) -> dict:
    """
    Converte il payload strutturato GTM nel formato standard con 'placeholders'.
    Supporta sia il formato diretto sia il payload strutturato.
    """
    placeholders: dict = {
        "customer_name":    body.get("customer_name", ""),
        "document_date":    body.get("document_date", ""),
        "document_version": body.get("document_version", "1.0"),
    }

    # ICP
    icp = body.get("icp", {})
    placeholders.update({
        "icp_industry":       icp.get("industry", ""),
        "icp_geography":      icp.get("geography", ""),
        "icp_company_size":   icp.get("company_size", ""),
        "icp_buying_group":   icp.get("buying_group", ""),
        "icp_maturity_level": icp.get("maturity_level", ""),
        "icp_custom_field":   icp.get("custom_field", ""),
    })

    # Persona
    p = body.get("persona", {})
    placeholders.update({
        "persona_role_name": p.get("role_name", ""),
        "persona_name":      p.get("name", ""),
        "persona_role":      p.get("role", ""),
        "persona_age":       p.get("age", ""),
        "persona_education": p.get("education", ""),
        "persona_location":  p.get("location", ""),
        "persona_industry":  p.get("industry", ""),
        "persona_channels":  p.get("channels", ""),
    })
    placeholders.update(_expand_list(p.get("responsibilities", []), "persona_resp", 3))

    for i in range(1, 5):
        obj = p.get(f"objective_{i}", {})
        pfx = f"persona_obj{i}"
        placeholders.update({
            f"{pfx}_title":   obj.get("title", f"Objective {i}"),
            f"{pfx}_desc":    obj.get("description", ""),
            f"{pfx}_kpi1":    obj.get("kpi1", ""),
            f"{pfx}_kpi2":    obj.get("kpi2", ""),
            f"{pfx}_kpi3":    obj.get("kpi3", ""),
            f"{pfx}_pain1":   obj.get("pain_point_1", ""),
            f"{pfx}_pain2":   obj.get("pain_point_2", ""),
            f"{pfx}_trigger": obj.get("winning_trigger", ""),
        })

    # Value Proposition
    vp = body.get("value_proposition", {})
    placeholders.update({
        "vp_what_is":   vp.get("what_is", ""),
        "vp_what_does": vp.get("what_does", ""),
        "vp_how_works": vp.get("how_works", ""),
        "vp_value":     vp.get("value_to_customer", ""),
    })

    # Slogans
    placeholders.update(_expand_list(body.get("slogans", []), "slogan", 3))

    # Customer Mental Model
    cmm = body.get("customer_mental_model", {})
    placeholders["cmm_persona"]   = cmm.get("persona", "")
    placeholders["cmm_objective"] = cmm.get("objective", "")
    placeholders.update(_expand_list(cmm.get("driving_factors", []), "cmm_driving_factor", 2))
    placeholders.update(_expand_list(cmm.get("pain_points", []),     "cmm_pain_point",     2))
    placeholders.update(_expand_list(cmm.get("use_cases", []),       "cmm_use_case",        2))

    # USPs
    usps = body.get("usps", [])
    for i in range(1, 4):
        u = usps[i - 1] if i - 1 < len(usps) else {}
        pfx = f"usp_{i}"
        placeholders[f"{pfx}_title"] = u.get("title", f"USP {i}")
        placeholders.update(_expand_list(u.get("body", []), f"{pfx}_body", 3))

    # Commercial Insight
    ci = body.get("commercial_insight", {})
    placeholders.update({
        "ci_objective":  ci.get("objective", ""),
        "ci_pain_point": ci.get("pain_point", ""),
        "ci_use_case":   ci.get("use_case", ""),
        "ci_usp":        ci.get("usp", ""),
        "ci_title":      ci.get("title", "Commercial Insight"),
    })
    placeholders.update(_expand_list(ci.get("narratives", []), "ci_narrative", 3))

    return {
        "template_s3_key": body.get("template_s3_key", GTM_TEMPLATE_KEY),
        "output_s3_key":   body.get("output_s3_key", "output/gtm_report_compilato.pptx"),
        "placeholders":    placeholders,
    }


def send_alert(error_msg: str, context) -> None:
    if not ALERT_EMAIL or not ALERT_FROM_EMAIL:
        return
    try:
        ses.send_email(
            Source=ALERT_FROM_EMAIL,
            Destination={"ToAddresses": [ALERT_EMAIL]},
            Message={
                "Subject": {
                    "Data": f"[va-ppt-connector] GTM Lambda Error – {context.function_name}",
                    "Charset": "UTF-8",
                },
                "Body": {
                    "Text": {
                        "Data": (
                            f"Errore nella Lambda {context.function_name}.\n\n"
                            f"Request ID: {context.aws_request_id}\n\n"
                            f"Errore:\n{error_msg}"
                        ),
                        "Charset": "UTF-8",
                    }
                },
            },
        )
    except Exception:
        pass


def replace_text_in_paragraph(paragraph, placeholders: dict) -> bool:
    """Sostituisce i placeholder nel paragrafo. Restituisce True se è avvenuta almeno una sostituzione."""
    full_text = "".join(run.text for run in paragraph.runs)
    if not PLACEHOLDER_RE.search(full_text):
        return False
    for key, value in placeholders.items():
        full_text = full_text.replace("{{" + key + "}}", str(value))
    if len(paragraph.runs) == 1:
        paragraph.runs[0].text = full_text
    elif paragraph.runs:
        paragraph.runs[0].text = full_text
        for run in paragraph.runs[1:]:
            run.text = ""
    return True


def _enable_autofit(shape) -> None:
    """Attiva il font auto-shrink sulla text frame (TEXT_TO_FIT_SHAPE via XML).

    python-pptx non espone MSO_AUTO_SIZE direttamente su tutte le shape,
    quindi operiamo direttamente sull'elemento XML <a:bodyPr>.
    """
    from lxml import etree
    txBody = shape.text_frame._txBody
    bodyPr = txBody.find("{http://schemas.openxmlformats.org/drawingml/2006/main}bodyPr")
    if bodyPr is None:
        return
    ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
    # Rimuovi eventuali figli di auto-size già presenti
    for tag in ("spAutoFit", "noAutofit", "normAutofit"):
        el = bodyPr.find(f"{{{ns}}}{tag}")
        if el is not None:
            bodyPr.remove(el)
    # Inserisci normAutofit (shrink font to fit shape)
    norm = etree.SubElement(bodyPr, f"{{{ns}}}normAutofit")
    # fontScale 100% (1000000 = 100%), lnSpcReduction 0 – PowerPoint gestirà il ridimensionamento
    norm.set("fontScale", "100000")
    norm.set("lnSpcReduction", "0")
    # Assicura word wrap
    bodyPr.set("wrap", "square")


def replace_placeholders(prs: Presentation, placeholders: dict) -> None:
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            replaced = False
            for paragraph in shape.text_frame.paragraphs:
                if replace_text_in_paragraph(paragraph, placeholders):
                    replaced = True
            if replaced:
                _enable_autofit(shape)


def lambda_handler(event, context):
    try:
        # --- Verifica API key ---
        if not check_api_key(event):
            return {
                "statusCode": 401,
                "headers": {"Content-Type": "application/json"},
                "body": json.dumps({"error": "Unauthorized: x-api-key mancante o non valida"}),
            }

        # --- Parsing input ---
        if isinstance(event.get("body"), str):
            body = json.loads(event["body"])
        else:
            body = event

        if isinstance(body, list):
            body = body[0]

        if "placeholders" not in body:
            body = preprocess_gtm(body)

        template_key = body["template_s3_key"]
        output_key   = body.get("output_s3_key", f"output/gtm_{context.aws_request_id}.pptx")
        placeholders = body.get("placeholders", {})

        # --- Download template da S3 ---
        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp_in:
            s3.download_fileobj(S3_BUCKET, template_key, tmp_in)
            tmp_in_path = tmp_in.name

        # --- Compila il template ---
        prs = Presentation(tmp_in_path)
        replace_placeholders(prs, placeholders)

        # --- Salva il risultato ---
        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp_out:
            prs.save(tmp_out.name)
            tmp_out_path = tmp_out.name

        # --- Upload su S3 ---
        with open(tmp_out_path, "rb") as f:
            s3.upload_fileobj(
                f, S3_BUCKET, output_key,
                ExtraArgs={"ContentType": "application/vnd.openxmlformats-officedocument.presentationml.presentation"},
            )

        os.unlink(tmp_in_path)
        os.unlink(tmp_out_path)

        # --- Pre-signed URL ---
        presigned_url = s3.generate_presigned_url(
            "get_object",
            Params={"Bucket": S3_BUCKET, "Key": output_key},
            ExpiresIn=PRESIGNED_URL_EXPIRATION,
        )

        return {
            "statusCode": 200,
            "headers": {"Content-Type": "application/json"},
            "body": json.dumps({
                "message": "GTM Report compilato con successo",
                "output_s3_key": output_key,
                "download_url": presigned_url,
            }),
        }

    except Exception:
        error_msg = traceback.format_exc()
        send_alert(error_msg, context)
        return {
            "statusCode": 500,
            "headers": {"Content-Type": "application/json"},
            "body": json.dumps({"error": error_msg}),
        }
