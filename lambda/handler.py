"""
AWS Lambda handler – compila un template PPTX sostituendo i placeholder
con i valori ricevuti nel JSON di input, salva il risultato su S3
e restituisce un pre-signed URL per il download.

Supporta due formati di input:
  1. Formato diretto:  {"template_s3_key": "...", "placeholders": {...}}
  2. Formato n8n:      [{"output": {...}}]  oppure  {"output": {...}}
                       con metadati opzionali nel body (customer_name, ecc.)
"""

import json
import os
import re
import tempfile
import traceback

import boto3
from pptx import Presentation

S3_BUCKET                = os.environ["S3_BUCKET"]
PRESIGNED_URL_EXPIRATION = int(os.environ.get("PRESIGNED_URL_EXPIRATION", "3600"))
ALERT_EMAIL              = os.environ.get("ALERT_EMAIL", "")
ALERT_FROM_EMAIL         = os.environ.get("ALERT_FROM_EMAIL", "")

s3  = boto3.client("s3")
ses = boto3.client("ses")

PLACEHOLDER_RE = re.compile(r"\{\{(\w+)\}\}")

# ---------------------------------------------------------------------------
# Mappatura: chiave nel JSON n8n → (prefisso placeholder template, max slot)
# ---------------------------------------------------------------------------
ARRAY_MAPPINGS = {
    "strengths_synthesis_items":     ("swot_strength",    4),
    "weaknesses_synthesis_items":    ("swot_weakness",    4),
    "opportunities_synthesis_items": ("swot_opportunity", 4),
    "threats_synthesis_items":       ("swot_threat",      4),
    "sales_so_finding":              ("sales_so_finding",     9),
    "sales_tb_finding":              ("sales_tb_finding",     3),
    "marketing_so_finding":          ("marketing_so_finding", 9),
    "marketing_tb_finding":          ("marketing_tb_finding", 3),
    "product_so_finding":            ("product_so_finding",   9),
    "product_tb_finding":            ("product_tb_finding",   3),
    # nel JSON n8n è "strategic_", nel template è "strategy_"
    "strategic_so_finding":          ("strategy_so_finding",  9),
    "strategic_tb_finding":          ("strategy_tb_finding",  3),
}


def _expand_array(items: list, prefix: str, max_items: int) -> dict:
    """Espande una lista in placeholder numerati: prefix_1 … prefix_N.
    Gli slot vuoti vengono impostati a stringa vuota."""
    return {
        f"{prefix}_{i}": items[i - 1] if i <= len(items) else ""
        for i in range(1, max_items + 1)
    }


def preprocess_n8n(body: dict) -> dict:
    """
    Converte il payload n8n nel formato standard atteso dalla Lambda.

    Il body può contenere:
      - "output": dict con gli array di analisi (obbligatorio per il formato n8n)
      - Metadati opzionali: customer_name, document_date, document_version,
        assessment_period, total_document_reviewed, total_interviews
      - "template_s3_key" e "output_s3_key" (opzionali, usano i default se assenti)
    """
    output_data = body.get("output", body)

    placeholders = {
        "customer_name":           body.get("customer_name", ""),
        "document_date":           body.get("document_date", ""),
        "document_version":        body.get("document_version", "1.0"),
        "assessment_period":       body.get("assessment_period", ""),
        "total_document_reviewed": str(body.get("total_document_reviewed", "")),
        "total_interviews":        str(body.get("total_interviews", "")),
    }

    for json_key, (template_prefix, max_n) in ARRAY_MAPPINGS.items():
        items = output_data.get(json_key, [])
        if isinstance(items, list):
            placeholders.update(_expand_array(items, template_prefix, max_n))

    return {
        "template_s3_key": body.get("template_s3_key", "templates/report_template.pptx"),
        "output_s3_key":   body.get("output_s3_key",   "output/report_compilato.pptx"),
        "placeholders":    placeholders,
    }


def send_alert(error_msg: str, context) -> None:
    """Invia un'email di alert via SES quando si verifica un errore."""
    if not ALERT_EMAIL or not ALERT_FROM_EMAIL:
        return
    try:
        ses.send_email(
            Source=ALERT_FROM_EMAIL,
            Destination={"ToAddresses": [ALERT_EMAIL]},
            Message={
                "Subject": {
                    "Data": f"[va-ppt-connector] Errore Lambda – {context.function_name}",
                    "Charset": "UTF-8",
                },
                "Body": {
                    "Text": {
                        "Data": (
                            f"Si è verificato un errore nella Lambda {context.function_name}.\n\n"
                            f"Request ID: {context.aws_request_id}\n\n"
                            f"Errore:\n{error_msg}"
                        ),
                        "Charset": "UTF-8",
                    }
                },
            },
        )
    except Exception:
        pass  # non bloccare il flusso per un errore nell'alert


def replace_text_in_paragraph(paragraph, placeholders: dict):
    """Sostituisce i placeholder {{key}} nel testo di un paragrafo.

    python-pptx spezza il testo in più *run* per motivi di formattazione,
    quindi ricostruiamo il testo completo, facciamo la sostituzione e poi
    lo rimappiamo nei run originali preservando lo stile.
    """
    full_text = "".join(run.text for run in paragraph.runs)
    if not PLACEHOLDER_RE.search(full_text):
        return

    for key, value in placeholders.items():
        full_text = full_text.replace("{{" + key + "}}", str(value))

    if len(paragraph.runs) == 1:
        paragraph.runs[0].text = full_text
    elif paragraph.runs:
        paragraph.runs[0].text = full_text
        for run in paragraph.runs[1:]:
            run.text = ""


def replace_placeholders(prs: Presentation, placeholders: dict):
    """Scorre tutte le slide e sostituisce i placeholder."""
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                replace_text_in_paragraph(paragraph, placeholders)


def lambda_handler(event, context):
    try:
        # --- Parsing input ---
        if isinstance(event.get("body"), str):
            body = json.loads(event["body"])
        else:
            body = event

        # Supporto array n8n: [{"output": {...}, "customer_name": "..."}]
        if isinstance(body, list):
            body = body[0]

        # Rilevamento formato: se manca "placeholders" è formato n8n → preprocess
        if "placeholders" not in body:
            body = preprocess_n8n(body)

        template_key = body["template_s3_key"]
        output_key   = body.get("output_s3_key", f"output/{context.aws_request_id}.pptx")
        placeholders = body.get("placeholders", {})

        # --- Download template da S3 ---
        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp_in:
            s3.download_fileobj(S3_BUCKET, template_key, tmp_in)
            tmp_in_path = tmp_in.name

        # --- Compila il template ---
        prs = Presentation(tmp_in_path)
        replace_placeholders(prs, placeholders)

        # --- Salva il risultato ---
        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp_out:
            prs.save(tmp_out.name)
            tmp_out_path = tmp_out.name

        # --- Upload su S3 ---
        with open(tmp_out_path, "rb") as f:
            s3.upload_fileobj(f, S3_BUCKET, output_key, ExtraArgs={
                "ContentType": "application/vnd.openxmlformats-officedocument.presentationml.presentation"
            })

        os.unlink(tmp_in_path)
        os.unlink(tmp_out_path)

        # --- Pre-signed URL ---
        presigned_url = s3.generate_presigned_url(
            "get_object",
            Params={"Bucket": S3_BUCKET, "Key": output_key},
            ExpiresIn=PRESIGNED_URL_EXPIRATION,
        )

        return {
            "statusCode": 200,
            "body": json.dumps({
                "message": "PPT compilato con successo",
                "output_s3_key": output_key,
                "download_url": presigned_url,
            }),
        }

    except Exception:
        error_msg = traceback.format_exc()
        send_alert(error_msg, context)
        return {
            "statusCode": 500,
            "body": json.dumps({"error": error_msg}),
        }
